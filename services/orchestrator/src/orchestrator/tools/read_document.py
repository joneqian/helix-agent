"""``read_document`` builtin — Tier 1 base capability.

A text-only main model (GLM / DeepSeek …) cannot natively read a PDF / Word /
Excel / PowerPoint a user dropped in the workspace. ``read_document`` parses
the common office + text formats to plain text so *any* agent can consume
them — no per-manifest opt-in, no code-execution detour. See
docs/design/agent-base-capabilities-and-form.md.

**Execution locus** — like the TE-7 file primitives it rides the agent's warm
sandbox ``exec`` channel (:func:`run_in_sandbox`): the parse runs *inside* the
per-user sandbox (not the control plane), so a malicious document's parser
exploit / zip-bomb is contained by the same gVisor + resource boundary as
every other tool. The parse libraries (pdfplumber / python-docx / openpyxl /
python-pptx) ship in the sandbox image (infra/sandbox-image/requirements.txt).

**Untrusted output** — a parsed document is attacker-influenceable content;
its text is spotlighted (datamarked + nonce-fenced) automatically by the
tools node like every other tool result (Stream PI-1b), so embedded
instructions read as data. The tool itself needs no special handling.

The snippet is *not* stdlib-only (unlike :mod:`file_ops`): it lazily imports a
parser per format, so an unsupported format or a missing parser returns a
clean ``{"ok": False, ...}`` envelope rather than crashing the sandbox.
"""

from __future__ import annotations

from collections.abc import Mapping
from dataclasses import dataclass
from typing import Any

from orchestrator.tools.file_ops import (
    _raise_for_error,
    _require_path,
    _snippet,
    parse_envelope,
)
from orchestrator.tools.registry import ToolContext, ToolResult, ToolSpec
from orchestrator.tools.sandbox import (
    SupervisorClient,
    run_in_sandbox,
)

_WORKSPACE_ROOT = "/workspace"
#: Largest document the parse will pull in. Documents are zips (docx/xlsx/pptx)
#: or page trees (pdf), so an oversized file is a decompression-bomb / OOM
#: risk — reject before parsing rather than mid-stream.
_MAX_DOC_BYTES = 25 * 1024 * 1024
#: Page / slide ceiling — a pathological 10k-page PDF must not blow the
#: sandbox memory or the agent's context window.
_MAX_DOC_PAGES = 200
#: Per-sheet row ceiling for spreadsheets (same rationale).
_MAX_SHEET_ROWS = 5000
#: Returned-text cap (chars). Documents are large; this bounds what reaches
#: the model's context. The full char count is reported in ``meta`` so the
#: model knows truncation happened.
_DOC_OUTPUT_CHAR_CAP = 200_000


# In-sandbox parse body. ``os`` / ``json`` / ``_P`` / ``_resolve`` come from
# the shared ``_PRELUDE`` (see file_ops). Parser imports are lazy + per-format
# so a missing library degrades to ``parser_unavailable`` not a crash.
_READ_DOCUMENT_MAIN = """

_TEXT_EXTS = {".txt", ".md", ".markdown", ".csv", ".tsv", ".log", ".json", ".rst", ".xml", ".html"}


def _extract(full, ext):
    if ext in _TEXT_EXTS:
        with open(full, "rb") as fh:
            return fh.read().decode("utf-8", "replace")
    if ext == ".pdf":
        import pdfplumber

        out = []
        with pdfplumber.open(full) as pdf:
            for page in pdf.pages[: _P["max_pages"]]:
                out.append(page.extract_text() or "")
        return "\\n\\n".join(out)
    if ext == ".docx":
        import docx

        doc = docx.Document(full)
        parts = [p.text for p in doc.paragraphs]
        for table in doc.tables:
            for row in table.rows:
                parts.append("\\t".join(cell.text for cell in row.cells))
        return "\\n".join(parts)
    if ext == ".xlsx":
        import openpyxl

        wb = openpyxl.load_workbook(full, read_only=True, data_only=True)
        try:
            out = []
            for ws in wb.worksheets:
                out.append("# " + ws.title)
                for i, row in enumerate(ws.iter_rows(values_only=True)):
                    if i >= _P["max_rows"]:
                        out.append("... (rows truncated)")
                        break
                    out.append("\\t".join("" if c is None else str(c) for c in row))
        finally:
            wb.close()
        return "\\n".join(out)
    if ext == ".pptx":
        import pptx

        prs = pptx.Presentation(full)
        out = []
        for n, slide in enumerate(list(prs.slides)[: _P["max_pages"]], 1):
            out.append("# Slide " + str(n))
            for shape in slide.shapes:
                if shape.has_text_frame:
                    out.append(shape.text_frame.text)
        return "\\n".join(out)
    return None


def _main():
    full = _resolve(_P["rel"])
    if full is None:
        return {"ok": False, "error": "path_escapes_workspace"}
    try:
        size = os.path.getsize(full)
    except FileNotFoundError:
        return {"ok": False, "error": "not_found"}
    except OSError as exc:
        return {"ok": False, "error": "io_error", "detail": str(exc)}
    if size > _P["max_bytes"]:
        return {"ok": False, "error": "file_too_large", "size": size}
    ext = os.path.splitext(full)[1].lower()
    try:
        text = _extract(full, ext)
    except ImportError:
        return {"ok": False, "error": "parser_unavailable", "format": ext.lstrip(".")}
    except Exception as exc:
        return {"ok": False, "error": "parse_failed", "detail": type(exc).__name__}
    if text is None:
        return {"ok": False, "error": "unsupported_format", "format": ext.lstrip(".")}
    cap = _P["cap"]
    return {
        "ok": True,
        "content": text[:cap],
        "format": ext.lstrip("."),
        "chars": len(text),
        "truncated": len(text) > cap,
    }


print(json.dumps(_main()))
"""


def build_read_document_wrapper(
    rel: str,
    *,
    cap: int,
    ws: str = _WORKSPACE_ROOT,
    max_bytes: int = _MAX_DOC_BYTES,
    max_pages: int = _MAX_DOC_PAGES,
    max_rows: int = _MAX_SHEET_ROWS,
) -> str:
    """Snippet that parses the document ``ws/rel`` and prints a JSON envelope."""
    return _snippet(
        {
            "ws": ws,
            "rel": rel,
            "cap": cap,
            "max_bytes": max_bytes,
            "max_pages": max_pages,
            "max_rows": max_rows,
        },
        _READ_DOCUMENT_MAIN,
    )


@dataclass
class ReadDocumentTool:
    """Parse a PDF / Word / Excel / PowerPoint / text file to plain text
    (exposed as ``read_document``)."""

    client: SupervisorClient
    output_char_cap: int = _DOC_OUTPUT_CHAR_CAP
    #: skill-runtime §5.1 — activated skill files seeded under /workspace/skills/.
    skill_seed_files: tuple[tuple[str, bytes], ...] = ()

    @property
    def spec(self) -> ToolSpec:
        return ToolSpec(
            name="read_document",
            description=(
                "Extract the text of a document in the agent's workspace — PDF, "
                "Word (.docx), Excel (.xlsx), PowerPoint (.pptx), or a plain-text "
                "format (.txt/.md/.csv/.json/...). Use this for binary documents "
                "that read_file cannot decode. Path is relative to /workspace."
            ),
            parameters={
                "type": "object",
                "properties": {
                    "path": {
                        "type": "string",
                        "description": "Workspace-relative document path (no leading '/' or '..').",
                    },
                },
                "required": ["path"],
            },
            is_read_only=True,
            path_args=("path",),
            side_effect="read_only",
            idempotent=True,
        )

    async def call(self, args: Mapping[str, Any], *, ctx: ToolContext) -> ToolResult:
        rel = _require_path(args, tool="read_document")
        outcome = await run_in_sandbox(
            self.client,
            code=build_read_document_wrapper(rel, cap=self.output_char_cap),
            timeout_s=None,
            ctx=ctx,
            tool_label="read_document",
            fallback_thread_id="read_document",
            seed_files=self.skill_seed_files,
        )
        env = parse_envelope(outcome, tool="read_document")
        _raise_for_error(env, tool="read_document")
        return ToolResult(
            content=str(env.get("content", "")),
            meta={
                "path": rel,
                "format": env.get("format"),
                "chars": env.get("chars"),
                "truncated": bool(env.get("truncated")),
            },
        )

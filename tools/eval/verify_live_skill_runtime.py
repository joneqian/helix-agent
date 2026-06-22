"""Live verification for skill runtime auto-mount + import classification.

Closes the "code green ≠ it actually runs" gap for skill-runtime §5.1/§5.2
(PRs #735-#738). Two phases against a running control-plane:

* **Phase 1 — import + classify (§5.2).** Import a skill from a public GitHub
  repo via the platform endpoint and assert the response carries a ``runtime``
  classification (``kind`` / ``runnable`` / ``hint``). system_admin only.

* **Phase 2 — auto-mount proof (§5.1).** Drive a *real* agent that has the
  imported skill bound + the ``bash`` tool, asking it to ``cat`` the seeded
  ``SKILL.md`` out of ``/workspace/skills/<name>/``. If the file is there, the
  supervisor (#736) + orchestrator (#737) seeding works end-to-end under the
  real sandbox runtime. This is the half CI can't do (no Docker/model in CI).

Keyless: the model key is resolved server-side; this script only sends prompts.
The API token (a **system_admin** dev-login bearer) is read from
``HELIX_API_TOKEN`` and never logged.

Usage (bring the dev stack up first — ``make dev-up``)::

    export HELIX_API_URL=http://localhost:8000    # control-plane (8080 is Keycloak)
    export HELIX_API_TOKEN=<system_admin dev bearer token>

    # Phase 1 only (no agent needed):
    uv run python tools/eval/verify_live_skill_runtime.py --import-only

    # One-shot end-to-end (import → activate skill → register the bundled
    # manifests/pptx-skill-test agent → auto-mount proof). Edit the manifest's
    # model block to your configured provider first:
    uv run python tools/eval/verify_live_skill_runtime.py --setup

    # Phase 2 against an agent you already registered (binds the skill + bash,
    # image_variant: office):
    uv run python tools/eval/verify_live_skill_runtime.py --agent pptx-agent@1.0.0

    # Phase 3 — office-image GENERATION (route ① §5.4). Requires the office
    # image built on the host (``make -C infra build-sandbox-office``) and a
    # model key configured. Registers the bundled manifests/pptx-office-gen-test
    # agent, runs a gated exec_python that writes a .pptx and converts it to .pdf
    # via soffice, approves the run, and asserts a non-empty PDF was produced:
    uv run python tools/eval/verify_live_skill_runtime.py --generate

Defaults import ``anthropics/skills`` → ``skills/pptx``. Exit code is non-zero
when any enabled phase fails — so it can gate a manual release check.
"""

from __future__ import annotations

import argparse
import asyncio
import json
import os
import sys
from pathlib import Path
from typing import Any

import httpx


def _require_env(name: str) -> str:
    value = os.environ.get(name)
    if not value:
        raise SystemExit(f"{name} is not set — export it before running (see module docstring)")
    return value


def _unwrap(data: dict[str, Any]) -> dict[str, Any]:
    """Unwrap the ``{success, data, error}`` envelope (sessions/agents routes)."""
    if data.get("success") is False:
        err = data.get("error") or {}
        raise SystemExit(f"API error: {err.get('code')}: {err.get('message')}")
    inner = data.get("data")
    return inner if isinstance(inner, dict) else data


def _iter_messages(obj: Any) -> list[dict[str, Any]]:
    found: list[dict[str, Any]] = []
    if isinstance(obj, dict):
        msgs = obj.get("messages")
        if isinstance(msgs, list):
            found.extend(m for m in msgs if isinstance(m, dict))
        for v in obj.values():
            found.extend(_iter_messages(v))
    elif isinstance(obj, list):
        for v in obj:
            found.extend(_iter_messages(v))
    return found


def _message_type(msg: dict[str, Any]) -> str:
    return str(msg.get("type") or msg.get("role") or "")


def _content_text(msg: dict[str, Any]) -> str:
    content = msg.get("content")
    if isinstance(content, str):
        return content
    if isinstance(content, list):
        return "".join(b.get("text", "") for b in content if isinstance(b, dict))
    return ""


# ── Phase 1 — import + classify ──────────────────────────────────────────────


async def phase_import(
    client: httpx.AsyncClient, *, source: str, skill: str
) -> tuple[bool, str | None, str | None]:
    """Import the skill from GitHub; assert a runtime classification is returned.

    Returns ``(ok, stored_skill_name, skill_id)``. The platform skills routes
    are raw (NOT enveloped), so the response is read directly.
    """
    print(f"[phase 1] import {source} #{skill}")
    resp = await client.post(
        "/v1/platform/skills/import-from-github",
        json={"source": source, "skill": skill},
    )
    if resp.status_code not in (200, 201):
        print(f"  FAIL — import HTTP {resp.status_code}: {resp.text[:300]}")
        return False, None, None
    body = resp.json()
    skill_obj = body.get("skill") or {}
    name = skill_obj.get("name")
    skill_id = skill_obj.get("id")
    runtime = body.get("runtime")
    print(f"  imported name={name!r} created={body.get('created')}")
    if not isinstance(runtime, dict):
        print("  FAIL — response carries no `runtime` classification (§5.2 regressed)")
        return False, name, skill_id
    print(
        f"  runtime: kind={runtime.get('kind')} runnable={runtime.get('runnable')}\n"
        f"           hint={runtime.get('hint')}"
    )
    return True, name, skill_id


async def activate_platform_skill(client: httpx.AsyncClient, *, skill_id: str) -> bool:
    """PATCH the platform skill to ACTIVE (imported skills start DRAFT, and an
    agent can only bind an ACTIVE platform skill). Idempotent."""
    print(f"[setup] activate platform skill {skill_id}")
    resp = await client.patch(f"/v1/platform/skills/{skill_id}", json={"status": "active"})
    if resp.status_code not in (200, 201):
        print(f"  FAIL — activate HTTP {resp.status_code}: {resp.text[:200]}")
        return False
    print(f"  status now: {resp.json().get('status')}")
    return True


async def register_agent(
    client: httpx.AsyncClient, *, manifest_path: str
) -> tuple[str, str] | None:
    """Register the agent from a manifest YAML; return ``(name, version)``.

    Idempotent: a 409 (already registered) is treated as success. Returns None
    on a real failure.
    """
    import yaml

    text = Path(manifest_path).read_text()
    meta = (yaml.safe_load(text) or {}).get("metadata", {})
    name, version = meta.get("name"), str(meta.get("version"))
    print(f"[setup] register agent {name}@{version} from {manifest_path}")
    resp = await client.post("/v1/agents", json={"manifest_yaml": text})
    if resp.status_code == 409:
        print("  already registered — reusing")
        return name, version
    if resp.status_code not in (200, 201):
        print(f"  FAIL — register HTTP {resp.status_code}: {resp.text[:300]}")
        return None
    print("  registered")
    return name, version


# ── Phase 2 — auto-mount proof ───────────────────────────────────────────────


async def _create_session(client: httpx.AsyncClient, name: str, version: str) -> str:
    resp = await client.post("/v1/sessions", json={"agent_name": name, "agent_version": version})
    resp.raise_for_status()
    return str(_unwrap(resp.json())["thread_id"])


class _RunTrace:
    """What the SSE stream revealed — to tell apart 'no bash call' from a run error."""

    def __init__(self) -> None:
        self.tool_msgs: list[tuple[str, str]] = []  # (tool_name, text)
        self.assistant_text: str = ""
        self.errors: list[str] = []
        self.events: dict[str, int] = {}
        self.raw_updates: list[str] = []


async def _consume_sse(resp: httpx.Response, tr: _RunTrace) -> None:
    """Parse an open SSE response into ``tr`` (event tally, tool/assistant msgs,
    errors). Shared by the run trigger and the resume continuation streams."""
    event = ""
    async for line in resp.aiter_lines():
        if line.startswith("event: "):
            event = line[len("event: ") :]
            tr.events[event] = tr.events.get(event, 0) + 1
        elif line.startswith("data: "):
            raw = line[len("data: ") :]
            if event in ("error", "run_error"):
                tr.errors.append(raw[:300])
                continue
            if event == "updates":
                tr.raw_updates.append(raw)
            try:
                payload = json.loads(raw)
            except json.JSONDecodeError:
                continue
            for msg in _iter_messages(payload):
                mtype = _message_type(msg)
                if mtype in ("tool", "ToolMessage"):
                    tr.tool_msgs.append((str(msg.get("name") or "?"), _content_text(msg)))
                elif mtype in ("ai", "AIMessage", "AIMessageChunk", "assistant"):
                    text = _content_text(msg)
                    if text.strip():
                        tr.assistant_text = text


async def _run_collect(
    client: httpx.AsyncClient, thread_id: str, prompt: str
) -> tuple[str | None, _RunTrace]:
    """Run the prompt; collect the trace + the run_id (from ``X-Helix-Run-Id``).

    The run_id matters for approval-gated runs: the SSE stream carries NO
    approval event — it just ends — so the caller polls the run by id to find
    the pause (see ``_wait_paused``)."""
    tr = _RunTrace()
    async with client.stream(
        "POST", f"/v1/sessions/{thread_id}/runs", json={"input": prompt}
    ) as resp:
        resp.raise_for_status()
        run_id = resp.headers.get("X-Helix-Run-Id")
        await _consume_sse(resp, tr)
    return run_id, tr


async def phase_mount(client: httpx.AsyncClient, *, agent: str, skill_name: str) -> bool:
    """Drive the agent to read the seeded SKILL.md; PASS if it's on disk."""
    name, _, version = agent.partition("@")
    if not version:
        raise SystemExit("--agent must be name@version")
    print(f"\n[phase 2] auto-mount proof via agent {agent}, skill={skill_name!r}")
    thread_id = await _create_session(client, name, version)
    print(f"  session: {thread_id}")

    # Probe via list_dir on the seeded skill dir (path is RELATIVE to /workspace
    # — no leading slash). A directory listing proves the files are on the
    # sandbox filesystem and is unfakeable by skill_view (which reads a single
    # file from the DB by exact path, and can't list a directory). list_dir is
    # read-only → NOT approval-gated (bash/exec_python are, side_effect=irreversible).
    rel_dir = f"skills/{skill_name}"
    prompt = (
        f"Use the list_dir tool to list the directory `{rel_dir}` (a path relative "
        "to the workspace root — do NOT prefix it with a slash). Report the entries "
        "exactly. Use ONLY list_dir; do not use skill_view or read_file."
    )
    _run_id, tr = await _run_collect(client, thread_id, prompt)
    print(f"  sse events: {tr.events}")

    # A real listing of the seeded dir contains SKILL.md alongside the bundled
    # scripts/ — skill_view of a single file never produces that pairing.
    def _is_listing(text: str) -> bool:
        return "[tool error]" not in text and "SKILL.md" in text and "scripts" in text

    if any(_is_listing(text) for _tool, text in tr.tool_msgs):
        print(f"  PASS — `{rel_dir}` listed off the sandbox filesystem (auto-mount works).")
        return True

    print("  FAIL — couldn't list the seeded skill dir off the sandbox filesystem.")
    print("  tool outputs:")
    for tool, text in tr.tool_msgs:
        print(f"    [{tool}] {text[:220].replace(chr(10), ' ')}")
    if not tr.tool_msgs:
        print("    <no tool calls>")
    if tr.assistant_text:
        print("  assistant said:")
        print("    " + tr.assistant_text[:300].replace("\n", "\n    "))
    print(
        "  NOTE: `not_found` / empty listing → files weren't seeded (is the running "
        "sandbox-supervisor rebuilt with the seed code? `make dev-up` now rebuilds it). "
        "`runner closed the connection` → sandbox didn't launch (image/runtime issue)."
    )
    return False


# ── Phase 3 — office-image generation through a gated run (route ① §5.4) ──────

# Verbatim code the agent is told to run via exec_python. Proves the office
# image's route-① binaries end-to-end: python-pptx writes a deck, then
# `soffice --headless` (the binary route ① added) converts it to PDF.
_GEN_CODE = """\
import os, subprocess
from pptx import Presentation

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "route 1 office gen 中文标题"
prs.save("/workspace/out.pptx")

proc = subprocess.run(
    ["soffice", "--headless", "--nolockcheck",
     "-env:UserInstallation=file:///workspace/.lo",
     "--convert-to", "pdf", "--outdir", "/workspace", "/workspace/out.pptx"],
    capture_output=True, text=True, timeout=120,
)
size = os.path.getsize("/workspace/out.pdf") if os.path.isfile("/workspace/out.pdf") else 0
print("GEN_RESULT", {"rc": proc.returncode, "pdf_exists": size > 0, "pdf_size": size})
"""


async def _wait_paused(
    client: httpx.AsyncClient, thread_id: str, run_id: str, *, attempts: int = 40
) -> dict[str, Any]:
    """Poll the run (raw JSON, not enveloped) until it pauses for approval or
    reaches a terminal state. Returns the last detail seen."""
    detail: dict[str, Any] = {}
    for _ in range(attempts):
        resp = await client.get(f"/v1/sessions/{thread_id}/runs/{run_id}")
        if resp.status_code == 200:
            detail = resp.json()
            status = detail.get("status")
            if status in ("paused", "completed", "succeeded", "failed", "cancelled"):
                return detail
        await asyncio.sleep(2)
    return detail


async def _resume_collect(client: httpx.AsyncClient, thread_id: str, run_id: str) -> _RunTrace:
    """Approve the paused run via the resume endpoint and collect the
    continuation stream (a fresh SSE stream of the resumed run)."""
    tr = _RunTrace()
    async with client.stream(
        "POST",
        f"/v1/sessions/{thread_id}/runs/{run_id}/resume",
        json={"decision": "approve"},
    ) as resp:
        resp.raise_for_status()
        ctype = resp.headers.get("content-type", "")
        if ctype.startswith("application/json"):
            # idempotent-replay branch (only if a key was reused — we don't send
            # one, so this is unexpected). Surface it rather than mis-parse.
            await resp.aread()
            tr.errors.append("resume returned JSON (idempotent replay?) instead of a stream")
            return tr
        await _consume_sse(resp, tr)
    return tr


async def phase_generate(client: httpx.AsyncClient, *, agent: str) -> bool:
    """Drive a real, approval-gated exec_python run that generates a .pptx and
    converts it to .pdf via soffice — the route-① office binaries, end-to-end."""
    name, _, version = agent.partition("@")
    if not version:
        raise SystemExit("--agent must be name@version")
    print(f"\n[phase 3] office-image generation via agent {agent}")
    thread_id = await _create_session(client, name, version)
    print(f"  session: {thread_id}")

    prompt = (
        "Use the exec_python tool to run EXACTLY this Python code, verbatim "
        "(do not modify it), then report the GEN_RESULT line it prints:\n"
        f"```python\n{_GEN_CODE}```"
    )
    run_id, _tr = await _run_collect(client, thread_id, prompt)
    if not run_id:
        print("  FAIL — no X-Helix-Run-Id header on the run (cannot track approval).")
        return False
    print(f"  run: {run_id} (streamed to end; approval is out-of-band)")

    detail = await _wait_paused(client, thread_id, run_id)
    status = detail.get("status")
    pending = detail.get("pending_approval")
    if status != "paused" or not pending:
        print(f"  FAIL — run did not pause for approval (status={status!r}).")
        print(
            "  NOTE: exec_python must be in the manifest's "
            "policies.approval_required_tools, else it never gates."
        )
        if status in ("failed", "cancelled"):
            print(f"  run detail: {json.dumps(detail)[:400]}")
        return False
    print(f"  paused on: {pending.get('action_summary')!r} → approving")

    tr = await _resume_collect(client, thread_id, run_id)
    print(f"  continuation sse events: {tr.events}")

    def _is_gen_ok(text: str) -> bool:
        # Tool stdout is wrapped in the spotlight/UNTRUSTED fence, which rewrites
        # spaces to ``▁`` — so a naive ``"'pdf_exists': True" in text`` misses the
        # real ``'pdf_exists':▁True``. Compare with all spacing stripped.
        flat = "".join(ch for ch in text if not ch.isspace() and ch != "▁")
        return "GEN_RESULT" in flat and "'pdf_exists':True" in flat

    if any(_is_gen_ok(text) for _tool, text in tr.tool_msgs):
        print("  PASS — agent generated .pptx and soffice produced .pdf in the office sandbox.")
        return True

    print("  FAIL — generation did not confirm a non-empty .pdf.")
    for tool, text in tr.tool_msgs:
        print(f"    [{tool}] {text[:300].replace(chr(10), ' ')}")
    if not tr.tool_msgs:
        print("    <no tool output in the continuation>")
    if tr.errors:
        print(f"  errors: {tr.errors}")
    print(
        "  NOTE: a soffice/poppler error or missing binary → the office image "
        "isn't rebuilt with route ① (run `make -C infra build-sandbox-office`)."
    )
    return False


# ── main ─────────────────────────────────────────────────────────────────────


_MANIFEST_DIR = Path(__file__).resolve().parents[2] / "manifests"
_DEFAULT_MANIFEST = _MANIFEST_DIR / "pptx-skill-test" / "v1.0.0.yaml"
_DEFAULT_GEN_MANIFEST = _MANIFEST_DIR / "pptx-office-gen-test" / "v1.0.0.yaml"


async def _amain(args: argparse.Namespace) -> int:
    base_url = args.base_url or _require_env("HELIX_API_URL")
    token = _require_env("HELIX_API_TOKEN")  # never logged
    headers = {"Authorization": f"Bearer {token}"}

    async with httpx.AsyncClient(base_url=base_url, headers=headers, timeout=180.0) as client:
        # Phase 3 (--generate) is standalone: it proves the office image binaries
        # via a gated exec_python run; no skill import/auto-mount involved.
        if args.generate:
            registered = await register_agent(client, manifest_path=args.gen_manifest)
            if registered is None:
                ok = False
            else:
                name, version = registered
                ok = await phase_generate(client, agent=f"{name}@{version}")
            print(f"\nRESULT: {'PASS — generation green.' if ok else 'FAIL — see above.'}")
            return 0 if ok else 1

        return await _run_import_phases(client, args)


async def _run_import_phases(client: httpx.AsyncClient, args: argparse.Namespace) -> int:
    """Phases 1-2: import + classify (§5.2) and auto-mount proof (§5.1)."""
    import_ok, skill_name, skill_id = await phase_import(
        client, source=args.source, skill=args.skill
    )
    ok = import_ok
    target_name = args.skill_name or skill_name

    if args.setup:
        # One-shot: activate the imported skill → register the agent →
        # auto-mount proof. Each step idempotent.
        if skill_id:
            ok = ok and await activate_platform_skill(client, skill_id=skill_id)
        registered = await register_agent(client, manifest_path=args.manifest)
        if registered is None:
            ok = False
        else:
            name, version = registered
            if not target_name:
                raise SystemExit("could not determine the stored skill name")
            ok = ok and await phase_mount(client, agent=f"{name}@{version}", skill_name=target_name)
    elif not args.import_only:
        if not args.agent:
            raise SystemExit("phase 2 needs --agent name@version (or --import-only / --setup)")
        if not target_name:
            raise SystemExit("could not determine the stored skill name; pass --skill-name")
        ok = ok and await phase_mount(client, agent=args.agent, skill_name=target_name)

    print(f"\nRESULT: {'PASS — all phases green.' if ok else 'FAIL — see above.'}")
    return 0 if ok else 1


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Live skill runtime auto-mount verification.")
    parser.add_argument("--base-url", default=None, help="control-plane URL (or $HELIX_API_URL)")
    parser.add_argument("--source", default="anthropics/skills", help="GitHub source (owner/repo)")
    parser.add_argument("--skill", default="skills/pptx", help="skill folder path in the repo")
    parser.add_argument("--agent", default=None, help="agent name@version for phase 2")
    parser.add_argument(
        "--skill-name",
        default=None,
        help="stored skill name for the /workspace path (else taken from import response)",
    )
    parser.add_argument("--import-only", action="store_true", help="run phase 1 only")
    parser.add_argument(
        "--setup",
        action="store_true",
        help="one-shot: import → activate skill → register agent → auto-mount proof",
    )
    parser.add_argument(
        "--manifest",
        default=str(_DEFAULT_MANIFEST),
        help="agent manifest YAML for --setup (default: manifests/pptx-skill-test/v1.0.0.yaml)",
    )
    parser.add_argument(
        "--generate",
        action="store_true",
        help="phase 3 (route ① §5.4): register the office-gen agent and prove a "
        "gated exec_python run generates a .pptx + soffice .pdf in the office image",
    )
    parser.add_argument(
        "--gen-manifest",
        default=str(_DEFAULT_GEN_MANIFEST),
        help="agent manifest YAML for --generate "
        "(default: manifests/pptx-office-gen-test/v1.0.0.yaml)",
    )
    args = parser.parse_args(argv)
    return asyncio.run(_amain(args))


if __name__ == "__main__":
    sys.exit(main())

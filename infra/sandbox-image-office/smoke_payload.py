"""Office-image smoke payload — runs INSIDE the sandbox via the runner.

Sent as the ``code`` of one runner request (``python -I -c``) by
``smoke_test.py``. Exercises the OFFICE-1b acceptance (design § 67): every
office library imports, a CJK xlsx/docx/pptx + a matplotlib chart generate
without error, and a Noto CJK font is resolvable (no tofu / missing-glyph).

Exit 0 + a trailing ``OK`` line == pass. Any failure raises, so the child
process exits non-zero and the runner reports ``exit_code != 0``.

matplotlib picks the headless Agg backend automatically here (no display in
the sandbox), so no explicit backend selection is needed.
"""

from __future__ import annotations

import importlib.util
import shutil
import subprocess
from pathlib import Path

import defusedxml
import imageio
import pandas as pd
import pdfplumber
import pypdf
from docx import Document
from matplotlib import font_manager
from matplotlib import pyplot as plt
from openpyxl import Workbook
from pdf2image import convert_from_path
from PIL import Image
from pptx import Presentation

_CJK = "中文办公能力验证 — 报表/图表"

# pip must be gone at runtime (no install path); installed packages still
# import — assert the hardening claim the image makes.
if importlib.util.find_spec("pip") is not None:
    raise RuntimeError("pip is still present in the office image")

# The baked matplotlibrc (OFFICE-ADR-3) must default to a CJK family with NO
# per-call config — that is its whole point. Check the default *before* the
# explicit override below, so a malformed/ignored baked rc fails the smoke.
_baked_default = plt.rcParams["font.sans-serif"]
if not any("CJK" in family for family in _baked_default):
    raise RuntimeError(f"baked matplotlibrc lost its CJK default: {_baked_default}")

# A CJK sans font must be registered, else matplotlib renders tofu for
# Chinese. Pick it dynamically (the .ttc face matplotlib registers may be
# named "...CJK JP"/"...CJK SC" depending on face order) rather than pinning
# one regional name. The image's baked matplotlibrc already defaults to these
# families; we assert one resolves and use it explicitly for the chart.
_cjk_names = sorted(
    n for n in {f.name for f in font_manager.fontManager.ttflist} if "Sans" in n and "CJK" in n
)
if not _cjk_names:
    raise RuntimeError("no Noto Sans CJK face registered by matplotlib")
cjk_font = _cjk_names[0]
plt.rcParams["font.sans-serif"] = [cjk_font, "DejaVu Sans"]
plt.rcParams["axes.unicode_minus"] = False

# xlsx (openpyxl) with Chinese cells.
wb = Workbook()
ws = wb.active
ws["A1"] = _CJK
ws["A2"] = 123
wb.save("/workspace/smoke.xlsx")

# docx with a Chinese paragraph.
doc = Document()
doc.add_paragraph(_CJK)
doc.save("/workspace/smoke.docx")

# pptx with a Chinese title slide.
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = _CJK
prs.save("/workspace/smoke.pptx")

# pandas round-trip through the xlsx engine.
df = pd.DataFrame({"名称": ["甲", "乙"], "值": [1, 2]})
df.to_excel("/workspace/df.xlsx", index=False)
if pd.read_excel("/workspace/df.xlsx").shape != (2, 2):
    raise RuntimeError("pandas xlsx round-trip changed shape")

# matplotlib chart with a Chinese title — the real CJK-render check.
fig, ax = plt.subplots()
ax.bar(["甲", "乙"], [1, 2])
ax.set_title(_CJK)
fig.savefig("/workspace/chart.png")

# Pillow can open the chart we just wrote.
with Image.open("/workspace/chart.png") as img:
    if img.size[0] <= 0:
        raise RuntimeError("Pillow read a zero-width chart")

# import is enough to prove these wheels + deps load (no separate exercise).
_ = (pypdf.__name__, pdfplumber.__name__, imageio.__name__, defusedxml.__name__)

# route ① §5.4 — the system binaries the office skills shell out to. Having
# the Python lib is not enough; assert the binary is present AND does real work.
for _bin in ("soffice", "pdftoppm", "ffmpeg"):
    if shutil.which(_bin) is None:
        raise RuntimeError(f"office binary missing from PATH: {_bin}")

# Real conversion: pptx → pdf via LibreOffice headless (the xlsx-recalc /
# pptx-thumbnail / docx-accept-changes paths all run soffice this way). A
# throwaway profile under the writable HOME avoids a stale-lock hang.
_soffice = subprocess.run(
    [  # noqa: S607 (PATH lookup of soffice is intentional)
        "soffice",
        "--headless",
        "--nolockcheck",
        "-env:UserInstallation=file:///workspace/.lo",
        "--convert-to",
        "pdf",
        "--outdir",
        "/workspace",
        "/workspace/smoke.pptx",
    ],
    capture_output=True,
    text=True,
    timeout=120,
    check=False,
)
_pdf = Path("/workspace/smoke.pdf")
if _soffice.returncode != 0 or not _pdf.is_file() or _pdf.stat().st_size == 0:
    raise RuntimeError(
        f"soffice pptx→pdf failed rc={_soffice.returncode} stderr={_soffice.stderr!r}"
    )

# Real PDF→image through pdf2image → proves the poppler binaries are wired.
_pages = convert_from_path(str(_pdf), dpi=50)
if not _pages or _pages[0].size[0] <= 0:
    raise RuntimeError("pdf2image/poppler produced no page image")

print(f"font={cjk_font}")
print("OK")
